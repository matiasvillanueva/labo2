"""Genera presentación PowerPoint del EDA de train_processed.csv."""

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).parent
DATA_PATH = BASE_DIR / "train_processed.csv"
CHARTS_DIR = BASE_DIR / "presentation_charts"
OUTPUT_PATH = BASE_DIR / "EDA_train_processed.pptx"

DARK = RGBColor(30, 58, 95)
ACCENT = RGBColor(0, 102, 153)
TEXT = RGBColor(40, 40, 40)
LIGHT = RGBColor(245, 247, 250)


def load_data():
    return pd.read_csv(DATA_PATH)


def generate_charts(df: pd.DataFrame) -> dict[str, Path]:
    CHARTS_DIR.mkdir(exist_ok=True)
    sns.set_theme(style="whitegrid", palette="muted")
    paths = {}

    # 1. Distribución log_price
    fig, ax = plt.subplots(figsize=(8, 4.5))
    sns.histplot(df["log_price"], kde=True, ax=ax, color="#336699")
    ax.set_title("Distribución de log_price", fontsize=13, fontweight="bold")
    ax.set_xlabel("log_price")
    p = CHARTS_DIR / "01_log_price_dist.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["price_dist"] = p

    # 2. Missing values
    missing = (df.isnull().sum() / len(df) * 100).sort_values(ascending=True)
    missing = missing[missing > 0].tail(10)
    fig, ax = plt.subplots(figsize=(8, 4.5))
    missing.plot(kind="barh", ax=ax, color="#cc6666")
    ax.set_xlabel("% faltante")
    ax.set_title("Top 10 columnas con valores faltantes", fontsize=13, fontweight="bold")
    p = CHARTS_DIR / "02_missing.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["missing"] = p

    # 3. homeType counts
    counts = df["homeType"].value_counts()
    fig, ax = plt.subplots(figsize=(8, 4.5))
    counts.plot(kind="barh", ax=ax, color="#336699")
    ax.set_xlabel("Cantidad de propiedades")
    ax.set_title("Distribución por tipo de propiedad", fontsize=13, fontweight="bold")
    p = CHARTS_DIR / "03_hometype.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["hometype"] = p

    # 4. Correlaciones con log_price
    corr = (
        df.select_dtypes(include=np.number)
        .corr()["log_price"]
        .drop(["log_price", "lastSoldPrice_hpi_adjusted"])
        .abs()
        .sort_values(ascending=True)
        .tail(10)
    )
    fig, ax = plt.subplots(figsize=(8, 4.5))
    corr.plot(kind="barh", ax=ax, color="#2d8f6f")
    ax.set_xlabel("|correlación| con log_price")
    ax.set_title("Top 10 correlaciones con el precio", fontsize=13, fontweight="bold")
    p = CHARTS_DIR / "04_correlations.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["corr"] = p

    # 5. Amenities impact
    amenities = ["has_pool", "has_garage", "has_waterfront", "has_hoa"]
    labels = ["Piscina", "Garaje", "Waterfront", "HOA"]
    diffs = []
    for col in amenities:
        m0 = df.loc[df[col] == 0, "log_price"].mean()
        m1 = df.loc[df[col] == 1, "log_price"].mean()
        diffs.append(m1 - m0)
    fig, ax = plt.subplots(figsize=(8, 4.5))
    colors = ["#2d8f6f" if d > 0 else "#cc6666" for d in diffs]
    ax.barh(labels, diffs, color=colors)
    ax.axvline(0, color="gray", linewidth=0.8)
    ax.set_xlabel("Diferencia en log_price (con vs sin)")
    ax.set_title("Impacto de amenities en el precio", fontsize=13, fontweight="bold")
    p = CHARTS_DIR / "05_amenities.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["amenities"] = p

    # 6. Mapa geográfico
    fig, ax = plt.subplots(figsize=(8, 5))
    sc = ax.scatter(
        df["longitude"],
        df["latitude"],
        c=df["log_price"],
        cmap="viridis",
        alpha=0.4,
        s=6,
    )
    plt.colorbar(sc, ax=ax, label="log_price")
    ax.set_xlabel("Longitud")
    ax.set_ylabel("Latitud")
    ax.set_title("Ubicación de propiedades (Florida)", fontsize=13, fontweight="bold")
    p = CHARTS_DIR / "06_geo_map.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["geo"] = p

    return paths


def set_slide_bg(slide, color: RGBColor = LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 220, 240)
    p2.alignment = PP_ALIGN.LEFT


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    image_path: Path | None = None,
    image_right: bool = True,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    text_width = Inches(5.8) if image_path else Inches(9)
    text_left = Inches(0.5)
    body_box = slide.shapes.add_textbox(text_left, Inches(1.1), text_width, Inches(5.8))
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
        p.level = 0

    if image_path and image_path.exists():
        img_left = Inches(6.5) if image_right else Inches(0.5)
        slide.shapes.add_picture(str(image_path), img_left, Inches(1.3), width=Inches(3.2))


def build_presentation(df: pd.DataFrame, charts: dict[str, Path]):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    price_mean = df["lastSoldPrice_hpi_adjusted"].mean()
    price_median = df["lastSoldPrice_hpi_adjusted"].median()

    # Slide 1 — Portada
    add_title_slide(
        prs,
        "EDA — Propiedades Inmobiliarias en Florida",
        "Análisis exploratorio de train_processed.csv\n11.840 propiedades · 46 variables",
    )

    # Slide 2 — Objetivo
    add_content_slide(
        prs,
        "Objetivo del análisis",
        [
            "Comprender la estructura y calidad del dataset de propiedades inmobiliarias.",
            "Identificar la variable objetivo y las features más relevantes para predecir precios.",
            "Detectar valores faltantes, patrones geográficos y diferencias entre tipos de propiedad.",
            "Documentar hallazgos que orienten el preprocesamiento y el modelado predictivo.",
            "El dataset proviene de listados de Zillow en Florida, con precios ajustados por HPI.",
        ],
    )

    # Slide 3 — Dataset overview
    add_content_slide(
        prs,
        "Descripción general del dataset",
        [
            "11.840 observaciones y 46 columnas (~9 MB en memoria).",
            "Cada fila representa una propiedad única identificada por zpid (ID de Zillow).",
            "Cobertura geográfica: Florida, con 64 zipcodes distintos.",
            "Tipos de variables: numéricas (35), categóricas (2) y texto libre (description).",
            "Incluye features originales, derivadas (logs, ratios, edad) y extraídas de texto.",
            "Target principal: log_price — logaritmo del precio de venta ajustado por HPI.",
        ],
    )

    # Slide 4 — Variables target e ID
    add_content_slide(
        prs,
        "Variables de identificación y target",
        [
            "zpid: identificador único de Zillow. No se usa como feature predictiva.",
            "lastSoldPrice_hpi_adjusted: precio de última venta en USD, ajustado por el Home Price Index para comparar ventas en distintos momentos.",
            "log_price: ln(precio ajustado). Variable objetivo del modelo.",
            f"Precio medio ajustado: USD {price_mean:,.0f} | Mediana: USD {price_median:,.0f}.",
            f"log_price: media = {df['log_price'].mean():.2f}, desvío = {df['log_price'].std():.2f}.",
            "La distribución de log_price es aproximadamente normal, lo que facilita modelos lineales.",
        ],
        charts["price_dist"],
    )

    # Slide 5 — Estructura física
    add_content_slide(
        prs,
        "Variables de estructura física",
        [
            "bedrooms: dormitorios (media ≈ 2.4). Faltante en ~4% de casos.",
            "bathrooms: baños, incluye medios baños (media ≈ 2.0).",
            "livingArea: superficie habitable en pies² (media ≈ 1.417 sqft).",
            "yearBuilt / property_age: año de construcción y edad derivada.",
            "lotAreaValue / log_lot_area: tamaño del lote. ~45% faltante (condos sin terreno propio).",
            "bath_to_bed_ratio y log_living_area: features derivadas para capturar no linealidades.",
            "homeType: categoría (SINGLE_FAMILY, CONDO, TOWNHOUSE, etc.) — fuerte separación de precios.",
        ],
    )

    # Slide 6 — Ubicación, impuestos, escuelas
    add_content_slide(
        prs,
        "Ubicación, impuestos y escuelas",
        [
            "latitude / longitude: coordenadas en Florida. Solo 31 valores faltantes (<0.3%).",
            "zipcode (5 dígitos) y zip_3digit: capturan efectos de mercado local.",
            "taxAssessedValue / latest_tax_value: valor tasado fiscal (~4% missing). Correlación ~0.70 con precio.",
            "latest_tax_paid / propertyTaxRate: impuesto pagado y tasa impositiva del condado.",
            "num_tax_records, num_sales, num_price_changes: historial de la propiedad.",
            "avg_school_rating / max_school_rating: calidad escolar cercana (correlación ~0.28–0.29).",
            "min_school_distance: proximidad a escuelas en millas.",
        ],
    )

    # Slide 7 — Amenities, tags, texto
    add_content_slide(
        prs,
        "Amenities, tags y texto",
        [
            "has_hoa / hoa_fee_monthly: asociación de propietarios y cuota mensual.",
            "has_pool, has_garage, has_waterfront: amenities binarias con impacto en precio.",
            "tag_price_cut, tag_new_construction, tag_foreclosure: etiquetas del listado.",
            "description: texto libre del aviso (11.804 descripciones únicas).",
            "desc_length / desc_word_count: longitud del texto.",
            "desc_is_boilerplate: indica descripciones autogeneradas (~mayoría del dataset).",
            "desc_mentions_*: flags si el texto menciona renovación, piscina, vista o 'nuevo'.",
        ],
    )

    # Slide 8 — Distribución del target
    add_content_slide(
        prs,
        "Análisis de la variable objetivo",
        [
            "log_price presenta distribución cercana a la normal, sin colas extremas.",
            "Rango: 10.84 a 14.50 (precios de ~USD 51K a ~USD 2M en escala original).",
            "Outliers por IQR: menos del 5% de las observaciones.",
            "El logaritmo estabiliza la varianza y reduce el efecto de propiedades muy caras.",
            "lastSoldPrice_hpi_adjusted tiene correlación 0.92 con log_price (es su transformación).",
            "La mediana (USD 472K) es menor que la media (USD 559K): ligera asimetría hacia precios altos.",
        ],
        charts["price_dist"],
    )

    # Slide 9 — Missing values
    missing_top = (df.isnull().sum() / len(df) * 100).sort_values(ascending=False)
    add_content_slide(
        prs,
        "Valores faltantes",
        [
            "15 columnas tienen al menos un valor faltante.",
            f"lotAreaValue y log_lot_area: {missing_top['lotAreaValue']:.1f}% — esperable en condos/apartments.",
            f"last_listing_price: {missing_top['last_listing_price']:.1f}% — propiedades sin historial de listado.",
            "taxAssessedValue, latest_tax_value: ~4.4% faltante.",
            "bedrooms, bathrooms, livingArea: <5% faltante — imputación sencilla.",
            "latitude/longitude: prácticamente completos (<0.3%).",
            "Estrategia sugerida: imputar área/lote por homeType; excluir o imputar listing price según modelo.",
        ],
        charts["missing"],
    )

    # Slide 10 — Tipos de propiedad
    ht = df["homeType"].value_counts()
    add_content_slide(
        prs,
        "Tipos de propiedad",
        [
            f"SINGLE_FAMILY: {ht['SINGLE_FAMILY']:,} ({ht['SINGLE_FAMILY']/len(df)*100:.0f}%) — segmento principal.",
            f"CONDO: {ht['CONDO']:,} ({ht['CONDO']/len(df)*100:.0f}%) — segundo grupo más grande.",
            f"TOWNHOUSE: {ht['TOWNHOUSE']:,} | APARTMENT: {ht['APARTMENT']:,} | MULTI_FAMILY: {ht['MULTI_FAMILY']:,}.",
            "SINGLE_FAMILY y TOWNHOUSE tienden a precios más altos que CONDO/APARTMENT.",
            "LOT y MANUFACTURED son categorías minoritarias (<1%).",
            "homeType es una variable clave: explica gran parte de la variabilidad estructural.",
            "Recomendación: usar encoding ordinal o one-hot según el modelo elegido.",
        ],
        charts["hometype"],
    )

    # Slide 11 — Correlaciones
    add_content_slide(
        prs,
        "Correlaciones con log_price",
        [
            "taxAssessedValue / latest_tax_value: ~0.70 — predictor más fuerte después del precio mismo.",
            "livingArea / log_living_area: ~0.48 — tamaño habitable muy relevante.",
            "last_listing_price: ~0.48 (con 33% missing — usar con cautela).",
            "bathrooms: ~0.40 | bedrooms: ~0.35 — estructura básica aporta señal.",
            "avg_school_rating / max_school_rating: ~0.28–0.29 — calidad escolar importa.",
            "latitude/longitude y zipcode capturan efectos geográficos no lineales.",
            "Riesgo: multicolinealidad entre tax variables y entre area/log_area.",
        ],
        charts["corr"],
    )

    # Slide 12 — Amenities
    add_content_slide(
        prs,
        "Impacto de amenities en el precio",
        [
            "has_garage: +0.34 en log_price (~+40% en escala original) — mayor impacto positivo.",
            "has_waterfront: +0.29 — frente al agua premium en Florida.",
            "has_pool: +0.18 — piscina incrementa el precio moderadamente.",
            "has_hoa: −0.11 — asociaciones pueden correlacionar con condos más económicos.",
            "tag_new_construction y tag_foreclosure: baja prevalencia pero señal de nicho.",
            "tag_price_cut: constante en todo el dataset (sin variación) — descartar como feature.",
            "Las amenities deben combinarse con homeType para evitar confusión de efectos.",
        ],
        charts["amenities"],
    )

    # Slide 13 — Geografía
    top_zips = df["zipcode"].value_counts().head(5)
    add_content_slide(
        prs,
        "Distribución geográfica",
        [
            "Propiedades concentradas en el sur de Florida (Miami-Dade, Broward, Palm Beach).",
            f"Zipcodes más frecuentes: {top_zips.index[0]}, {top_zips.index[1]}, {top_zips.index[2]} (~500 c/u).",
            "El mapa muestra gradientes de precio: costa este generalmente más cara.",
            "zip_3digit agrupa 16 zonas — útil para modelos con menos cardinalidad que zipcode.",
            "Coordenadas permiten modelos espaciales o clustering geográfico.",
            "La ubicación explica variación que las features estructurales no capturan solas.",
        ],
        charts["geo"],
    )

    # Slide 14 — Hallazgos
    add_content_slide(
        prs,
        "Hallazgos principales",
        [
            "Dataset robusto: 11.840 propiedades, target bien distribuido, pocas columnas críticas incompletas.",
            "Valor tasado fiscal es el mejor predictor individual (r ≈ 0.70).",
            "Tipo de propiedad y ubicación segmentan fuertemente el mercado.",
            "Condos explican la mayoría de missing en lotAreaValue — no es error de datos.",
            "Descripciones son mayormente boilerplate; features de texto tienen baja prevalencia.",
            "Multicolinealidad presente entre variables de impuestos y área — requiere selección de features.",
            "El dataset está listo para modelado con preprocesamiento moderado.",
        ],
    )

    # Slide 15 — Cierre
    add_title_slide(
        prs,
        "Conclusiones y próximos pasos",
        "1. Imputar missing values por tipo de propiedad\n"
        "2. Seleccionar features (eliminar redundancias)\n"
        "3. Encoding de categóricas (homeType, zipcode)\n"
        "4. Entrenar modelos baseline (Lineal, RF, XGBoost)\n"
        "5. Evaluar con validación cruzada y métricas (RMSE, MAE, R²)",
    )

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


def main():
    df = load_data()
    charts = generate_charts(df)
    output = build_presentation(df, charts)
    print(f"Presentación creada: {output}")
    print(f"Gráficos guardados en: {CHARTS_DIR}")


if __name__ == "__main__":
    main()
